import os
import shutil
from datetime import datetime
from mcp.server.fastmcp import FastMCP

# Try to import Office libraries (fail silently if not installed)
try:
    import docx
except ImportError:
    docx = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    import pptx
except ImportError:
    pptx = None

try:
    import PyPDF2
except ImportError:
    PyPDF2 = None

# Default workspace directory
DEFAULT_WORKSPACE = "D:\\ALICE"

# Default safe zones - can be configured
DEFAULT_SAFE_ZONES = [
    os.path.expanduser("~/Documents"),
    os.path.expanduser("~/Downloads"),
    os.path.expanduser("~/Desktop"),
]


def register_file_tools(mcp: FastMCP, safe_zones: list[str] = None, base_url: str = ""):
    """
    Register file management tools with tiered permissions.

    Permission Model:
    - READ operations: Allowed anywhere
    - WRITE operations: Only in safe zones
    - COPY: Can copy INTO safe zones from anywhere, but NOT out of safe zones

    Args:
        mcp: FastMCP instance
        safe_zones: List of directories where write operations are allowed.
        base_url: Base URL of the server (e.g., "http://localhost:9999")
    """
    allowed_zones = safe_zones or DEFAULT_SAFE_ZONES
    # Normalize paths
    allowed_zones = [os.path.abspath(z) for z in allowed_zones]

    # Ensure workspace exists
    os.makedirs(DEFAULT_WORKSPACE, exist_ok=True)

    def is_in_safe_zone(path: str) -> bool:
        """Check if path is within any safe zone."""
        try:
            abs_path = os.path.abspath(path)
            # Normalize path case for Windows
            norm_path = os.path.normcase(abs_path)

            for zone in allowed_zones:
                abs_zone = os.path.abspath(zone)
                norm_zone = os.path.normcase(abs_zone)

                # Check if it's the zone itself
                if norm_path == norm_zone:
                    return True

                # Check if it's a file inside the zone
                # Add separator to ensure we don't match partial folder names
                # e.g. C:\Users\User\Doc shouldn't match C:\Users\User\Documents
                if not norm_zone.endswith(os.sep):
                    norm_zone += os.sep

                if norm_path.startswith(norm_zone):
                    return True

            return False
        except Exception:
            return False

    def get_safe_zones_str() -> str:
        """Get formatted list of safe zones for error messages."""
        return "\n".join(f"  - {z}" for z in allowed_zones)

    # ==================== HELPERS FOR OFFICE FILES ====================

    def _read_docx(path: str, max_lines: int) -> str:
        if not docx:
            return "Error: python-docx not installed. Cannot read .docx files."
        try:
            doc = docx.Document(path)
            full_text = []
            for para in doc.paragraphs:
                if len(full_text) >= max_lines:
                    full_text.append(f"\n... (truncated at {max_lines} paragraphs)")
                    break
                full_text.append(para.text)
            return "\n".join(full_text)
        except Exception as e:
            return f"Error reading DOCX: {str(e)}"

    def _read_xlsx(path: str, max_lines: int) -> str:
        if not openpyxl:
            return "Error: openpyxl not installed. Cannot read .xlsx files."
        try:
            wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
            result = []

            # Read active sheet only
            sheet = wb.active
            result.append(f"[Sheet: {sheet.title}]")

            row_count = 0
            for row in sheet.iter_rows(values_only=True):
                if row_count >= max_lines:
                    result.append(f"\n... (truncated at {max_lines} rows)")
                    break

                # Filter None values and join
                row_text = "\t".join([str(cell) if cell is not None else "" for cell in row])
                result.append(row_text)
                row_count += 1

            return "\n".join(result)
        except Exception as e:
            return f"Error reading XLSX: {str(e)}"

    def _read_pptx(path: str, max_lines: int) -> str:
        if not pptx:
            return "Error: python-pptx not installed. Cannot read .pptx files."
        try:
            prs = pptx.Presentation(path)
            text_runs = []

            slide_count = 0
            for slide in prs.slides:
                slide_count += 1
                text_runs.append(f"--- Slide {slide_count} ---")

                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)

                if len(text_runs) >= max_lines:
                    text_runs.append(f"\n... (truncated at {max_lines} text blocks)")
                    break

            return "\n".join(text_runs)
        except Exception as e:
            return f"Error reading PPTX: {str(e)}"

    def _read_pdf(path: str, max_lines: int) -> str:
        if not PyPDF2:
            return "Error: PyPDF2 not installed. Cannot read .pdf files."
        try:
            reader = PyPDF2.PdfReader(path)
            text = []

            count = 0
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text.append(f"--- Page {count + 1} ---")
                    text.append(page_text)
                    count += 1

                if count >= max_lines: # Treat pages as lines/blocks
                    text.append(f"\n... (truncated at {max_lines} pages)")
                    break

            return "\n".join(text)
        except Exception as e:
            return f"Error reading PDF: {str(e)}"

    # ==================== WORKSPACE INFO ====================

    @mcp.tool(name="MyPC-get_workspace")
    def get_workspace() -> str:
        """
        Get the default workspace directory path.

        Default workspace is D:\ALICE. All file operations can use this as the base directory.
        You can use relative paths or just the workspace for operations.

        Returns:
            Default workspace path and information.
        """
        return f"Default Workspace: {DEFAULT_WORKSPACE}\n\nThis is the base directory for all file operations."

    # ==================== READ OPERATIONS (Anywhere) ====================

    @mcp.tool(name="MyPC-list_directory")
    def list_directory(path: str = None) -> str:
        """
        List contents of a directory. (READ - allowed anywhere)

        Default workspace: D:\ALICE (use path="D:\\ALICE" or omit for workspace)

        Args:
            path: Absolute path to the directory. If omitted or empty, uses workspace (D:\\ALICE).

        Returns:
            Formatted list of files and folders with their sizes and modification times.
        """
        # Use workspace if path not provided
        if not path:
            path = DEFAULT_WORKSPACE

        if not os.path.exists(path):
            return f"Error: Path does not exist: {path}"

        if not os.path.isdir(path):
            return f"Error: Not a directory: {path}"

        try:
            entries = []
            for name in os.listdir(path):
                full_path = os.path.join(path, name)
                try:
                    stat = os.stat(full_path)
                    mtime = datetime.fromtimestamp(stat.st_mtime).strftime("%Y-%m-%d %H:%M")

                    if os.path.isdir(full_path):
                        entries.append(f"[DIR]  {name}/  ({mtime})")
                    else:
                        size = stat.st_size
                        if size < 1024:
                            size_str = f"{size} B"
                        elif size < 1024 * 1024:
                            size_str = f"{size // 1024} KB"
                        else:
                            size_str = f"{size // (1024 * 1024)} MB"
                        entries.append(f"[FILE] {name}  ({size_str}, {mtime})")
                except PermissionError:
                    entries.append(f"[???]  {name}  (Permission Denied)")

            if not entries:
                return f"Directory is empty: {path}"

            return f"Contents of {path}:\n" + "\n".join(sorted(entries))

        except PermissionError:
            return f"Error: Permission denied to access: {path}"
        except Exception as e:
            return f"Error listing directory: {str(e)}"

    @mcp.tool(name="MyPC-read_file")
    def read_file(path: str, max_lines: int = 500) -> str:
        """
        Read contents of a file. Supports Text, Word (.docx), Excel (.xlsx), PPT (.pptx), PDF.
        (READ - allowed anywhere)

        Default workspace: D:\ALICE

        Args:
            path: Absolute path to the file. Can be in workspace (D:\\ALICE\\file.txt).
            max_lines: Maximum number of lines/blocks to read (default 500).

        Returns:
            File contents as text.
        """
        if not os.path.exists(path):
            return f"Error: File does not exist: {path}"

        if not os.path.isfile(path):
            return f"Error: Not a file: {path}"

        ext = os.path.splitext(path)[1].lower()

        # Route to appropriate handler
        if ext == ".docx":
            return _read_docx(path, max_lines)
        elif ext == ".xlsx":
            return _read_xlsx(path, max_lines)
        elif ext == ".pptx":
            return _read_pptx(path, max_lines)
        elif ext == ".pdf":
            return _read_pdf(path, max_lines)

        # Default text reading for other files
        try:
            with open(path, 'r', encoding='utf-8', errors='replace') as f:
                lines = []
                for i, line in enumerate(f):
                    if i >= max_lines:
                        lines.append(f"\n... (truncated at {max_lines} lines)")
                        break
                    lines.append(line.rstrip('\n\r'))

            return "\n".join(lines)

        except PermissionError:
            return f"Error: Permission denied to read: {path}"
        except Exception as e:
            return f"Error reading file: {str(e)}"

    @mcp.tool(name="MyPC-get_file_info")
    def get_file_info(path: str) -> str:
        """
        Get detailed information about a file or directory. (READ - allowed anywhere)

        Default workspace: D:\ALICE

        Args:
            path: Path to the file or directory. Can be in workspace (D:\\ALICE\\file.txt).

        Returns:
            Detailed file information.
        """
        if not os.path.exists(path):
            return f"Error: Path does not exist: {path}"

        try:
            stat = os.stat(path)
            abs_path = os.path.abspath(path)
            in_safe = "Yes" if is_in_safe_zone(path) else "No"

            info = [
                f"Path: {abs_path}",
                f"Type: {'Directory' if os.path.isdir(path) else 'File'}",
                f"Size: {stat.st_size} bytes",
                f"Created: {datetime.fromtimestamp(stat.st_ctime).strftime('%Y-%m-%d %H:%M:%S')}",
                f"Modified: {datetime.fromtimestamp(stat.st_mtime).strftime('%Y-%m-%d %H:%M:%S')}",
                f"Accessed: {datetime.fromtimestamp(stat.st_atime).strftime('%Y-%m-%d %H:%M:%S')}",
                f"In Safe Zone: {in_safe}",
            ]

            if os.path.isdir(path):
                try:
                    count = len(os.listdir(path))
                    info.append(f"Items: {count}")
                except PermissionError:
                    info.append("Items: (Permission Denied)")

            return "\n".join(info)

        except PermissionError:
            return f"Error: Permission denied to access: {path}"
        except Exception as e:
            return f"Error getting file info: {str(e)}"

    @mcp.tool(name="MyPC-list_safe_zones")
    def list_safe_zones() -> str:
        """
        List all configured safe zones where write operations are allowed.

        Returns:
            List of safe zone directories.
        """
        return "Safe Zones (write operations allowed):\n" + get_safe_zones_str()

    @mcp.tool(name="MyPC-get_download_url")
    def get_download_url(path: str) -> str:
        """
        Get a direct download URL for a file.

        The file MUST be located in a Safe Zone.

        Args:
            path: Absolute path to the file.

        Returns:
            Download URL if valid, or error message.
        """
        if not os.path.exists(path):
            return f"Error: File does not exist: {path}"

        if not os.path.isfile(path):
            return f"Error: Not a file: {path}"

        if not is_in_safe_zone(path):
            return f"Error: Download denied. File must be in a safe zone.\n\nSafe Zones:\n{get_safe_zones_str()}"

        try:
            # URL encode the path
            import urllib.parse
            encoded_path = urllib.parse.quote(path)
            url = f"{base_url}/download?path={encoded_path}"
            return f"Download URL: {url}"
        except Exception as e:
            return f"Error generating URL: {str(e)}"

    # ==================== WRITE OPERATIONS (Safe Zones Only) ====================

    @mcp.tool(name="MyPC-edit_file")
    def edit_file(path: str, old_text: str, new_text: str, count: int = 1) -> str:
        """
        Edit a file by replacing specific text. (WRITE - safe zones only)

        Performs an exact string replacement. Preserves the rest of the file content.

        Args:
            path: Absolute path to the file.
            old_text: The text to be replaced (must match exactly, including newlines/indentation).
            new_text: The new text to insert.
            count: Number of occurrences to replace (default 1). Use -1 to replace all.

        Returns:
            Success or error message.
        """
        if not is_in_safe_zone(path):
            return f"Error: Edit operation denied. Path must be in a safe zone.\n\nSafe Zones:\n{get_safe_zones_str()}"

        if not os.path.exists(path):
            return f"Error: File does not exist: {path}"

        try:
            # Read file
            with open(path, 'r', encoding='utf-8') as f:
                content = f.read()

            # Check if old_text exists
            if old_text not in content:
                # Provide a helpful hint if it's a whitespace issue
                normalized_content = " ".join(content.split())
                normalized_old = " ".join(old_text.split())
                if normalized_old in normalized_content:
                    return "Error: Text found but with different whitespace/indentation. Please provide exact match."
                return "Error: The 'old_text' was not found in the file."

            # Perform replacement
            new_content = content.replace(old_text, new_text, count)

            # Check if anything actually changed (in case count was 0 or other logic)
            if new_content == content:
                return "Warning: No changes were made (replacement logic returned same content)."

            # Write back
            with open(path, 'w', encoding='utf-8') as f:
                f.write(new_content)

            return f"File edited successfully: {path}"

        except Exception as e:
            return f"Error editing file: {str(e)}"
    def write_file(path: str, content: str) -> str:
        """
        Write content to a file. (WRITE - safe zones only)

        Default workspace: D:\ALICE (recommended for new files)
        Safe zones: Documents, Downloads, Desktop, Pictures, D:\\, E:\\

        Args:
            path: Absolute path to the file (must be in a safe zone).
                   Example: "D:\\ALICE\\notes.txt" or just "notes.txt" for workspace.
            content: Text content to write.

        Returns:
            Success or error message.
        """
        # If path is not absolute, use workspace
        if not os.path.isabs(path):
            path = os.path.join(DEFAULT_WORKSPACE, path)

        if not is_in_safe_zone(path):
            return f"Error: Write operation denied. Path must be in a safe zone.\n\nSafe Zones:\n{get_safe_zones_str()}"

        try:
            parent = os.path.dirname(path)
            if parent and not os.path.exists(parent):
                os.makedirs(parent, exist_ok=True)

            with open(path, 'w', encoding='utf-8') as f:
                f.write(content)

            return f"File written successfully: {path}"

        except Exception as e:
            return f"Error writing file: {str(e)}"

    @mcp.tool(name="MyPC-move_file")
    def move_file(source: str, destination: str) -> str:
        """
        Move or rename a file/directory. (WRITE - both paths must be in safe zones)

        Args:
            source: Source path (must be in safe zone).
            destination: Destination path (must be in safe zone).

        Returns:
            Success or error message.
        """
        if not is_in_safe_zone(source):
            return f"Error: Cannot move from outside safe zone.\nSource: {source}\n\nSafe Zones:\n{get_safe_zones_str()}"

        if not is_in_safe_zone(destination):
            return f"Error: Cannot move to outside safe zone.\nDestination: {destination}\n\nSafe Zones:\n{get_safe_zones_str()}"

        if not os.path.exists(source):
            return f"Error: Source does not exist: {source}"

        try:
            shutil.move(source, destination)
            return f"Moved: {source} -> {destination}"

        except Exception as e:
            return f"Error moving file: {str(e)}"

    @mcp.tool(name="MyPC-delete_file")
    def delete_file(path: str, to_recycle: bool = True) -> str:
        """
        Delete a file or directory. (WRITE - safe zones only)

        Args:
            path: Path to delete (must be in safe zone).
            to_recycle: If True, move to recycle bin. If False, permanently delete.

        Returns:
            Success or error message.
        """
        if not is_in_safe_zone(path):
            return f"Error: Delete operation denied. Path must be in a safe zone.\n\nSafe Zones:\n{get_safe_zones_str()}"

        if not os.path.exists(path):
            return f"Error: Path does not exist: {path}"

        try:
            if to_recycle:
                try:
                    from send2trash import send2trash
                    send2trash(path)
                    return f"Moved to Recycle Bin: {path}"
                except ImportError:
                    return "Error: send2trash not installed. Use to_recycle=False for permanent deletion."
            else:
                if os.path.isdir(path):
                    shutil.rmtree(path)
                else:
                    os.remove(path)
                return f"Permanently deleted: {path}"

        except Exception as e:
            return f"Error deleting: {str(e)}"

    @mcp.tool(name="MyPC-create_directory")
    def create_directory(path: str) -> str:
        """
        Create a new directory. (WRITE - safe zones only)

        Default workspace: D:\ALICE (recommended for new folders)

        Args:
            path: Path for the new directory (must be in a safe zone).
                   Example: "D:\\ALICE\\projects" or just "projects" for workspace.

        Returns:
            Success or error message.
        """
        # If path is not absolute, use workspace
        if not os.path.isabs(path):
            path = os.path.join(DEFAULT_WORKSPACE, path)

        if not is_in_safe_zone(path):
            return f"Error: Create directory denied. Path must be in a safe zone.\n\nSafe Zones:\n{get_safe_zones_str()}"

        try:
            os.makedirs(path, exist_ok=True)
            return f"Directory created: {path}"

        except Exception as e:
            return f"Error creating directory: {str(e)}"

    # ==================== COPY OPERATION (Special: INTO safe zone only) ====================

    @mcp.tool(name="MyPC-copy_file")
    def copy_file(source: str, destination: str) -> str:
        """
        Copy a file or directory. (SPECIAL - destination must be in safe zone)

        Security: You can copy FROM anywhere INTO a safe zone, but NOT out of safe zones.

        Args:
            source: Source path (can be anywhere).
            destination: Destination path (must be in a safe zone).

        Returns:
            Success or error message.
        """
        if not is_in_safe_zone(destination):
            return f"Error: Copy destination must be in a safe zone.\nDestination: {destination}\n\nSafe Zones:\n{get_safe_zones_str()}"

        if not os.path.exists(source):
            return f"Error: Source does not exist: {source}"

        try:
            if os.path.isdir(source):
                shutil.copytree(source, destination)
            else:
                # Ensure parent directory exists
                parent = os.path.dirname(destination)
                if parent and not os.path.exists(parent):
                    os.makedirs(parent, exist_ok=True)
                shutil.copy2(source, destination)

            return f"Copied: {source} -> {destination}"

        except Exception as e:
            return f"Error copying: {str(e)}"
